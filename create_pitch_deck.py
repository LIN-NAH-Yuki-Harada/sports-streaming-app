#!/usr/bin/env python3
"""Sports Streaming App 企画書 PPT生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === カラーパレット ===
BG_DARK = RGBColor(0x0D, 0x1B, 0x2A)       # ダークネイビー
BG_SECTION = RGBColor(0x13, 0x24, 0x3B)     # セクション背景
ACCENT = RGBColor(0x00, 0xC8, 0x53)         # グリーン（スポーツ感）
ACCENT2 = RGBColor(0x00, 0xA3, 0xFF)        # ブルー
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC4)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED_ACCENT = RGBColor(0xFF, 0x45, 0x45)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=LIGHT_GRAY, font_name="Meiryo"):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(6)
    return txBox


def add_card(slide, left, top, width, height, title, body_lines, title_color=ACCENT, bg_color=BG_SECTION):
    """カード型のコンポーネントを追加"""
    shape = add_shape_fill(slide, left, top, width, height, bg_color)
    # 角丸風の見た目（枠線で代用）
    shape.line.color.rgb = RGBColor(0x2A, 0x3A, 0x55)
    shape.line.width = Pt(1)

    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
                 title, font_size=20, color=title_color, bold=True)

    y_offset = top + Inches(0.7)
    for line in body_lines:
        text = line[0] if isinstance(line, tuple) else line
        color = line[1] if isinstance(line, tuple) and len(line) > 1 else LIGHT_GRAY
        size = line[2] if isinstance(line, tuple) and len(line) > 2 else 14
        add_text_box(slide, left + Inches(0.3), y_offset, width - Inches(0.6), Inches(0.35),
                     text, font_size=size, color=color)
        y_offset += Inches(0.32)
    return shape


# ============================================================
# Slide 1: 表紙
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

# アクセントライン
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

# メインタイトル
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Sports Streaming App", font_size=54, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "〜 すべての子どもの試合が、どこからでも見られる世界をつくる 〜",
             font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

# サブタイトル
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
             "ローカルスポーツ特化 UGC型ストリーミングプラットフォーム 企画書",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 日付
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "2026年4月  |  Confidential",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 下部アクセントライン
add_shape_fill(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), ACCENT)


# ============================================================
# Slide 2: 課題 — 現状の問題点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), RED_ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "PROBLEM", font_size=14, color=RED_ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "ローカルスポーツ配信の「3つの壁」", font_size=36, color=WHITE, bold=True)

# 課題カード1
add_card(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5),
         "1. 配信されない試合",
         [
             ("スポーツブル等の大手は", ),
             ("全国大会・プロの試合のみ", RED_ACCENT, 16),
             ("", ),
             ("地域の大会、部活の公式戦、", ),
             ("スポ少の試合は配信されず、", ),
             ("見に行けない家族は結果待ち", ),
         ], title_color=RED_ACCENT)

# 課題カード2
add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.5),
         "2. 既存配信の映像問題",
         [
             ("YouTube Live等で配信しても…", ),
             ("", ),
             ("得点板が見えない", RED_ACCENT, 16),
             ("情報表示で映像を遮ってしまう", RED_ACCENT, 16),
             ("「今のスコアは？」に応えられない", RED_ACCENT, 16),
             ("", ),
             ("→ 視聴者のストレスが溜まる",),
         ], title_color=RED_ACCENT)

# 課題カード3
add_card(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.5),
         "3. 配信者の負担",
         [
             ("撮影しながらスコアを", ),
             ("更新するのは困難", RED_ACCENT, 16),
             ("", ),
             ("1人で撮影 + 情報管理は", ),
             ("現実的に無理がある", ),
             ("", ),
             ("→ 配信の質が下がる or", ),
             ("   そもそも配信を諦める", ),
         ], title_color=RED_ACCENT)


# ============================================================
# Slide 3: 解決策 — プロダクト概要
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "SOLUTION", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "スマホだけでTV中継品質を実現する", font_size=36, color=WHITE, bold=True)

# コア機能 3本柱
add_card(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.8),
         "TV中継風オーバーレイ",
         [
             ("スコア・チーム名・時間を", ),
             ("映像に常時オーバーレイ表示", ACCENT, 15),
             ("", ),
             ("映像を遮ることなく", ),
             ("試合情報が常に視聴者に届く", ),
             ("", ),
             ("「今知りたい！」に", ),
             ("常に応え続けるUI", ACCENT, 15),
         ], title_color=ACCENT)

add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.8),
         "専用リモコン",
         [
             ("別のスマホがスコア入力", ),
             ("専用リモコンに変身", ACCENT, 15),
             ("", ),
             ("得点・イベント（ゴール、", ),
             ("ファウル等）をワンタップ入力", ),
             ("→ 即座にオーバーレイ反映", ),
             ("", ),
             ("撮影者は撮影に集中できる", ACCENT, 15),
         ], title_color=ACCENT)

add_card(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.8),
         "マルチカメラ AI切替",
         [
             ("複数スマホをWiFi/BT/", ),
             ("クラウドで紐付け", ),
             ("", ),
             ("AIが自動でベストアングルを", ),
             ("選択・カット切替", ACCENT, 15),
             ("", ),
             ("保護者2〜3人のスマホで", ),
             ("プロ品質の中継が完成", ACCENT, 15),
         ], title_color=ACCENT)


# ============================================================
# Slide 4: ユーザー体験フロー
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT2)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "USER EXPERIENCE", font_size=14, color=ACCENT2, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "配信者・視聴者それぞれの体験", font_size=36, color=WHITE, bold=True)

# 配信者サイド
add_shape_fill(slide, Inches(0.5), Inches(1.9), Inches(5.9), Inches(0.5), ACCENT)
add_text_box(slide, Inches(0.8), Inches(1.95), Inches(5), Inches(0.4),
             "配信者（保護者）の体験", font_size=16, color=DARK_TEXT, bold=True)

flow_lines_broadcaster = [
    "1.  アプリを開いて「配信開始」",
    "2.  試合情報（チーム名・競技）を入力",
    "3.  スマホで撮影開始 → 自動でオーバーレイ表示",
    "4.  別の保護者がリモコンアプリでスコア入力",
    "5.  さらに別の保護者がサブカメラとして参加",
    "6.  AIが自動でカメラ切替 → TV中継品質の完成！",
]
y = Inches(2.6)
for line in flow_lines_broadcaster:
    add_text_box(slide, Inches(0.8), y, Inches(5.5), Inches(0.35),
                 line, font_size=15, color=LIGHT_GRAY)
    y += Inches(0.38)

# 視聴者サイド
add_shape_fill(slide, Inches(6.9), Inches(1.9), Inches(5.9), Inches(0.5), ACCENT2)
add_text_box(slide, Inches(7.2), Inches(1.95), Inches(5), Inches(0.4),
             "視聴者（家族・関係者）の体験", font_size=16, color=WHITE, bold=True)

flow_lines_viewer = [
    "1.  アプリ or Webでログイン",
    "2.  お気に入りのチーム・大会をフォロー",
    "3.  ライブ配信の通知が届く",
    "4.  タップして視聴開始",
    "5.  スコア・試合状況がリアルタイムで見える",
    "6.  見逃しても → アーカイブでいつでも再視聴",
]
y = Inches(2.6)
for line in flow_lines_viewer:
    add_text_box(slide, Inches(7.2), y, Inches(5.5), Inches(0.35),
                 line, font_size=15, color=LIGHT_GRAY)
    y += Inches(0.38)

# ポイント
add_shape_fill(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2), BG_SECTION)
add_text_box(slide, Inches(1), Inches(5.75), Inches(11), Inches(0.9),
             "ポイント:  保護者2〜3人のスマホだけで、プロのTV中継チームと同等の配信が可能。\n"
             "特別な機材・技術知識は一切不要。アプリが全てをサポートする。",
             font_size=16, color=ACCENT)


# ============================================================
# Slide 5: ターゲット市場
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "TARGET MARKET", font_size=14, color=ORANGE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "膨大な未開拓市場", font_size=36, color=WHITE, bold=True)

# 市場データ
add_card(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(2.5),
         "小学校スポーツ少年団",
         [
             ("全国約 25,000 団", WHITE, 18),
             ("登録団員数 約 60万人", ),
             ("→ その背後に120万人の保護者", ),
         ], title_color=ORANGE)

add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(2.5),
         "中学校部活動",
         [
             ("全国約 10,000 校", WHITE, 18),
             ("運動部員数 約 200万人", ),
             ("地区大会・県大会は配信ゼロ", ),
         ], title_color=ORANGE)

add_card(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(2.5),
         "高校部活動",
         [
             ("全国約 5,000 校", WHITE, 18),
             ("インターハイ予選等の地方大会", ),
             ("注目度は高いが配信がない", ),
         ], title_color=ORANGE)

# まとめ
add_shape_fill(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), BG_SECTION)
add_multi_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(1.8), [
    ("対象スポーツ: 全ジャンル", 16, WHITE, True),
    ("サッカー / 野球 / バスケ / バレー / 陸上 / 柔道 / 剣道 / 水泳 / テニス / 卓球 … etc.", 14, LIGHT_GRAY),
    ("", 8),
    ("既存サービス（スポーツブル等）がカバーしていないローカル大会・地域大会は、", 15, LIGHT_GRAY),
    ("年間 数十万試合 規模の膨大な未開拓市場。", 18, ORANGE, True),
])


# ============================================================
# Slide 6: 収益モデル
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "BUSINESS MODEL", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "2段階の収益モデル", font_size=36, color=WHITE, bold=True)

# フェーズ1
add_card(slide, Inches(0.5), Inches(2.0), Inches(5.9), Inches(4.5),
         "Phase 1: サブスクリプション",
         [
             ("", ),
             ("月額課金で視聴し放題", WHITE, 18),
             ("", ),
             ("「有料でも絶対に価値がある」", ),
             ("と感じてもらえるサービス品質が前提", ),
             ("", ),
             ("想定プラン:", WHITE, 15),
             ("  ベーシック:  月額 ¥500〜（1チームフォロー）",),
             ("  スタンダード: 月額 ¥980〜（全試合見放題）",),
             ("  ファミリー:  月額 ¥1,480〜（家族共有）",),
         ], title_color=ACCENT)

# フェーズ2
add_card(slide, Inches(6.9), Inches(2.0), Inches(5.9), Inches(4.5),
         "Phase 2: 広告収入（ユーザー拡大後）",
         [
             ("", ),
             ("スポーツ親和性の高い企業の広告", WHITE, 18),
             ("", ),
             ("想定広告主:", WHITE, 15),
             ("  スポーツ用品メーカー（ミズノ、アシックス等）", ),
             ("  飲料メーカー（アクエリアス、ポカリ等）", ),
             ("  保険会社（スポーツ保険）", ),
             ("  学習塾・教育サービス", ),
             ("", ),
             ("マス層浸透で広告単価UP", ACCENT, 15),
         ], title_color=ACCENT)


# ============================================================
# Slide 7: Exit戦略
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT2)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "EXIT STRATEGY", font_size=14, color=ACCENT2, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "スポーツブルとのシナジー", font_size=36, color=WHITE, bold=True)

# 現状 vs 連携後
add_card(slide, Inches(0.5), Inches(2.0), Inches(5.9), Inches(2.2),
         "現状のスポーツブル",
         [
             ("自社カメラクルーで撮影 → コスト大", ),
             ("全国大会・プロのみ → コンテンツ量に限界", ),
             ("ローカル大会は完全に未対応", ),
         ], title_color=LIGHT_GRAY)

add_card(slide, Inches(6.9), Inches(2.0), Inches(5.9), Inches(2.2),
         "連携後のスポーツブル",
         [
             ("保護者UGC → 撮影コストゼロ", ACCENT, 15),
             ("ローカル大会が大量追加 → コンテンツ爆増", ACCENT, 15),
             ("UU数の大幅UP → 広告収益UP", ACCENT, 15),
         ], title_color=ACCENT)

# Win-Win
add_shape_fill(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), BG_SECTION)
add_multi_text(slide, Inches(1), Inches(4.8), Inches(11), Inches(2.2), [
    ("買収時の価値提案", 20, ACCENT2, True),
    ("", 8),
    ("1. UGC配信基盤  — 保護者が撮影するインフラが完成済み", 16, WHITE),
    ("2. アクティブユーザー  — ローカルスポーツの保護者コミュニティ", 16, WHITE),
    ("3. コンテンツ量  — 年間数十万試合分のアーカイブ", 16, WHITE),
    ("4. 統合コスト最小  — API連携を前提とした設計", 16, WHITE),
])


# ============================================================
# Slide 8: 開発ロードマップ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "ROADMAP", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "段階的開発ロードマップ", font_size=36, color=WHITE, bold=True)

# MVP
add_card(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.8),
         "MVP",
         [
             ("基本配信 + スコアオーバーレイ", WHITE, 16),
             ("", ),
             ("1台のスマホで配信", ),
             ("スコア手動入力", ),
             ("TV風オーバーレイ表示", ),
             ("基本的な視聴画面（Web）", ),
             ("ユーザー登録・認証", ),
             ("", ),
             ("目標: 動くプロトタイプ", ACCENT, 14),
             ("「これは使える」の検証", ACCENT, 14),
         ], title_color=ACCENT)

# v1.0
add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.8),
         "v1.0",
         [
             ("リモコン連携", WHITE, 16),
             ("", ),
             ("別デバイスからスコア操作", ),
             ("リアルタイム同期", ),
             ("チーム・大会管理機能", ),
             ("モバイルアプリ（iOS/Android）", ),
             ("サブスク課金機能", ),
             ("", ),
             ("目標: 正式リリース", ACCENT2, 14),
             ("有料ユーザー獲得開始", ACCENT2, 14),
         ], title_color=ACCENT2)

# v2.0
add_card(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.8),
         "v2.0",
         [
             ("マルチカメラ + AI切替", WHITE, 16),
             ("", ),
             ("複数スマホ連携", ),
             ("AI自動カメラスイッチング", ),
             ("広告配信基盤", ),
             ("スポーツブル連携API", ),
             ("アナリティクス・ダッシュボード", ),
             ("", ),
             ("目標: Exit準備完了", ORANGE, 14),
             ("買収提案可能な状態", ORANGE, 14),
         ], title_color=ORANGE)


# ============================================================
# Slide 9: 開発体制
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT2)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "TEAM & DEVELOPMENT", font_size=14, color=ACCENT2, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "開発体制", font_size=36, color=WHITE, bold=True)

# 体制
add_card(slide, Inches(0.5), Inches(2.0), Inches(5.9), Inches(2.5),
         "コアチーム",
         [
             ("オーナー — プロジェクト指揮・意思決定・プロダクト設計", WHITE, 15),
             ("Claude AI — 設計・開発・技術アドバイザー・コードレビュー", WHITE, 15),
             ("", ),
             ("AI駆動の少人数開発で、スピードと品質を両立", ACCENT2, 14),
         ], title_color=ACCENT2)

add_card(slide, Inches(6.9), Inches(2.0), Inches(5.9), Inches(2.5),
         "外部リソース（必要に応じて）",
         [
             ("ココナラ等で個人開発者に依頼", WHITE, 15),
             ("", ),
             ("想定発注領域:", LIGHT_GRAY, 14),
             ("  UI/UXデザイン、ネイティブアプリ実装、", ),
             ("  映像処理・ストリーミング技術、インフラ構築", ),
         ], title_color=ACCENT2)

# 開発方針
add_shape_fill(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.2), BG_SECTION)
add_multi_text(slide, Inches(1), Inches(5.1), Inches(11), Inches(2.0), [
    ("開発方針", 20, ACCENT2, True),
    ("", 8),
    ("1. MVP First  — まず動くものを作り、実際のユーザーで検証する", 15, WHITE),
    ("2. AI活用  — Claude によるコード生成・レビューで開発速度を最大化", 15, WHITE),
    ("3. 段階的外注  — コア機能はClaudeと内製、専門領域はココナラで外注", 15, WHITE),
    ("4. 技術選定は柔軟に  — 外注しやすい技術スタックを選ぶ", 15, WHITE),
])


# ============================================================
# Slide 10: 技術概要（外注用にも使える）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "TECH OVERVIEW", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "技術構成（想定）", font_size=36, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5),
         "フロントエンド",
         [
             ("Web: Next.js / React", WHITE, 15),
             ("Mobile: React Native or Flutter", WHITE, 15),
             ("", ),
             ("オーバーレイ描画: Canvas / WebGL", ),
             ("リアルタイム同期: WebSocket", ),
             ("", ),
             ("→ Web/Mobile共通の技術基盤で", ACCENT, 13),
             ("   開発効率を最大化", ACCENT, 13),
         ], title_color=ACCENT)

add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.5),
         "バックエンド",
         [
             ("API: Node.js or Go", WHITE, 15),
             ("DB: Supabase (PostgreSQL)", WHITE, 15),
             ("認証: Supabase Auth", WHITE, 15),
             ("", ),
             ("ストリーミング: WebRTC / HLS", ),
             ("リアルタイム: Supabase Realtime", ),
             ("", ),
             ("→ Supabase活用で", ACCENT, 13),
             ("   インフラ管理コスト最小化", ACCENT, 13),
         ], title_color=ACCENT)

add_card(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.5),
         "AI / インフラ",
         [
             ("AI切替: 映像解析 + ルールベース", WHITE, 15),
             ("デバイス間通信: WebRTC P2P", WHITE, 15),
             ("", ),
             ("クラウド: AWS or GCP", ),
             ("CDN: CloudFront / Cloudflare", ),
             ("CI/CD: GitHub Actions", ),
             ("", ),
             ("→ スケーラブルかつ", ACCENT, 13),
             ("   外注しやすい構成", ACCENT, 13),
         ], title_color=ACCENT)


# ============================================================
# Slide 11: ココナラ外注計画
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "OUTSOURCING PLAN", font_size=14, color=ORANGE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "外注戦略（ココナラ活用）", font_size=36, color=WHITE, bold=True)

# 内製 vs 外注
add_card(slide, Inches(0.5), Inches(2.0), Inches(5.9), Inches(2.0),
         "内製（オーナー + Claude AI）",
         [
             ("プロダクト設計・要件定義・プロジェクト管理", WHITE, 15),
             ("バックエンドAPI・DB設計・基本的なフロントエンド", WHITE, 15),
             ("ビジネスロジック・コードレビュー", WHITE, 15),
         ], title_color=ACCENT)

add_card(slide, Inches(6.9), Inches(2.0), Inches(5.9), Inches(2.0),
         "外注候補（ココナラ等）",
         [
             ("UI/UXデザイン（アプリ・Web）", WHITE, 15),
             ("ネイティブアプリ開発（iOS/Android固有機能）", WHITE, 15),
             ("映像処理・ストリーミング実装の専門家", WHITE, 15),
         ], title_color=ORANGE)

# 外注の進め方
add_shape_fill(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7), BG_SECTION)
add_multi_text(slide, Inches(1), Inches(4.6), Inches(11), Inches(2.5), [
    ("外注の進め方", 20, ORANGE, True),
    ("", 8),
    ("Step 1:  MVPを内製で完成させ、プロダクトの方向性を固める", 15, WHITE),
    ("Step 2:  MVP検証後、専門性が必要な領域を特定して発注", 15, WHITE),
    ("Step 3:  要件定義書・技術仕様書をClaude AIで作成し、発注精度を上げる", 15, WHITE),
    ("Step 4:  コードレビュー・品質管理はClaude AIでカバー", 15, WHITE),
    ("", 8),
    ("ポイント: 「何を作るか」が明確な状態で発注することで、コスト・品質をコントロール", 14, ORANGE),
])


# ============================================================
# Slide 12: まとめ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(1.0),
             "Sports Streaming App", font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
             "すべての子どもの試合が、どこからでも見られる世界をつくる",
             font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

# まとめポイント
summary_items = [
    ("未開拓市場", "スポーツブル等がカバーしない年間数十万試合のローカルスポーツ"),
    ("UGC × テクノロジー", "保護者のスマホだけでTV中継品質を実現する3つのコア機能"),
    ("明確なExit", "スポーツブルへの買収提案を見据えた設計・モデリング"),
    ("AI駆動の開発", "Claude AI + ココナラ外注で、少人数でもスピーディに開発"),
]

y = Inches(2.8)
for title, desc in summary_items:
    add_shape_fill(slide, Inches(2), y, Inches(9.3), Inches(0.8), BG_SECTION)
    add_text_box(slide, Inches(2.3), y + Inches(0.08), Inches(2.5), Inches(0.35),
                 title, font_size=18, color=ACCENT, bold=True)
    add_text_box(slide, Inches(2.3), y + Inches(0.4), Inches(8.5), Inches(0.35),
                 desc, font_size=14, color=LIGHT_GRAY)
    y += Inches(0.95)

# CTA
add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.6),
             "Next Step → MVP開発開始",
             font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# 保存
# ============================================================
output_path = "/Users/yukiharada/Desktop/Sports-Streeming/Sports_Streaming_企画書.pptx"
prs.save(output_path)
print(f"Done: {output_path}")
