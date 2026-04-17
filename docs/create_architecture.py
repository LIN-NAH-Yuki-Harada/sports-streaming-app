#!/usr/bin/env python3
"""LIVE SPOtCH システム構造チャート — シンプル版 + 将来拡張"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

FONT = "Meiryo UI"

BG      = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x0F, 0x17, 0x2A)
TEXT    = RGBColor(0x1A, 0x1A, 0x2E)
SUB     = RGBColor(0x66, 0x66, 0x80)
PRIMARY = RGBColor(0xE6, 0x39, 0x46)
P_LIGHT = RGBColor(0xFD, 0xEB, 0xED)
BLUE    = RGBColor(0x25, 0x63, 0xEB)
B_LIGHT = RGBColor(0xEB, 0xF0, 0xFF)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
G_LIGHT = RGBColor(0xE8, 0xF5, 0xEB)
ORANGE  = RGBColor(0xEA, 0x88, 0x0B)
O_LIGHT = RGBColor(0xFF, 0xF3, 0xE0)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
P2_LIGHT= RGBColor(0xF3, 0xEE, 0xFF)
GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
BORDER  = RGBColor(0xE0, 0xE0, 0xE8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_L  = RGBColor(0xF5, 0xF5, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = BG

def rect(l, t, w, h, color, lc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if lc: sh.line.color.rgb = lc; sh.line.width = Pt(1.5)
    else: sh.line.fill.background()
    return sh

def rrect(l, t, w, h, color, lc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if lc: sh.line.color.rgb = lc; sh.line.width = Pt(1.5)
    else: sh.line.fill.background()
    return sh

def txt(l, t, w, h, text, sz=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = FONT
    p.alignment = align
    return box

def arrow_d(l, t, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def arrow_r(l, t, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.NOTCHED_RIGHT_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def arrow_lr(l, t, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def line_v(l, t, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Pt(2), h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

# ═══ 左端アクセント ═══
rect(Inches(0), Inches(0), Pt(5), Inches(7.5), PRIMARY)

# ═══ タイトル ═══
txt(Inches(0.6), Inches(0.25), Inches(7), Inches(0.4),
    "LIVE SPOtCH — システム構造", sz=22, color=DARK, bold=True)
txt(Inches(0.6), Inches(0.65), Inches(7), Inches(0.25),
    "現在の構成と将来の拡張ポイント", sz=11, color=SUB)
rect(Inches(0.6), Inches(0.95), Inches(1), Pt(3), PRIMARY)

# ════════════════════════════════════════
# 左半分: 現在のシステム構造（紐付き図）
# ════════════════════════════════════════

# --- ユーザー ---
rrect(Inches(2.5), Inches(1.2), Inches(3.2), Inches(0.7), P_LIGHT, PRIMARY)
txt(Inches(2.7), Inches(1.28), Inches(2.8), Inches(0.25),
    "ユーザー（保護者・家族）", sz=13, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
txt(Inches(2.7), Inches(1.55), Inches(2.8), Inches(0.2),
    "スマホのブラウザからアクセス（アプリDL不要）", sz=8, color=SUB, align=PP_ALIGN.CENTER)

# 矢印: ユーザー → アプリ
arrow_d(Inches(3.9), Inches(1.95), Inches(0.35), Inches(0.35), GRAY)

# --- アプリ（中央の大きなボックス） ---
rrect(Inches(0.6), Inches(2.4), Inches(7.0), Inches(2.7), GRAY_L, DARK)
txt(Inches(0.9), Inches(2.5), Inches(6.4), Inches(0.3),
    "LIVE SPOtCH（Next.js / Vercel）", sz=14, color=DARK, bold=True)

# アプリ内の機能ブロック（横並び）
app_features = [
    ("配信エンジン", "カメラ映像取得\nスコアオーバーレイ\nリアルタイム配信", PRIMARY),
    ("視聴画面", "共有コードで接続\nスコア自動更新\nLINE共有", BLUE),
    ("チーム管理", "招待コード\nメンバー権限\nチーム配信", GREEN),
    ("決済", "Stripe連携\n¥300/¥500月額\n解約・退会", ORANGE),
    ("YouTube", "OAuth連携\nアーカイブ保存\n（Phase 2-3）", PURPLE),
]
for i, (title, desc, ac) in enumerate(app_features):
    x = Inches(0.8 + i * 1.36)
    rrect(x, Inches(2.95), Inches(1.18), Inches(1.85), WHITE, ac)
    rect(x, Inches(2.95), Inches(1.18), Pt(3), ac)
    txt(x + Inches(0.08), Inches(3.05), Inches(1.02), Inches(0.2),
        title, sz=10, color=ac, bold=True, align=PP_ALIGN.CENTER)
    txt(x + Inches(0.08), Inches(3.3), Inches(1.02), Inches(1.3),
        desc, sz=8, color=SUB, align=PP_ALIGN.CENTER)

# --- 矢印: アプリ → 外部サービス ---
# 各機能から下の対応サービスへの縦線
connections = [
    (Inches(1.39), PRIMARY),   # 配信 → LiveKit
    (Inches(2.75), BLUE),      # 視聴 → Supabase
    (Inches(4.11), GREEN),     # チーム → Supabase
    (Inches(5.47), ORANGE),    # 決済 → Stripe
    (Inches(6.50), PURPLE),    # YouTube → YouTube API
]
for x, color in connections:
    line_v(x, Inches(4.85), Inches(0.45), color)

# --- 外部サービス ---
y_svc = Inches(5.35)
services = [
    ("LiveKit Cloud", "映像配信基盤", PRIMARY, P_LIGHT),
    ("Supabase", "DB・認証・Realtime", BLUE, B_LIGHT),
    ("Stripe", "決済処理", ORANGE, O_LIGHT),
    ("YouTube API", "動画保存", PURPLE, P2_LIGHT),
    ("Vercel", "ホスティング", GREEN, G_LIGHT),
]
for i, (title, desc, ac, bg) in enumerate(services):
    x = Inches(0.6 + i * 1.42)
    rrect(x, y_svc, Inches(1.24), Inches(0.85), bg, ac)
    txt(x + Inches(0.06), y_svc + Inches(0.1), Inches(1.12), Inches(0.2),
        title, sz=10, color=ac, bold=True, align=PP_ALIGN.CENTER)
    txt(x + Inches(0.06), y_svc + Inches(0.38), Inches(1.12), Inches(0.3),
        desc, sz=8, color=SUB, align=PP_ALIGN.CENTER)

# Supabaseは2つの機能（視聴+チーム）に紐づくので横線を追加
rect(Inches(2.75), Inches(5.05), Inches(1.36), Pt(2), BLUE)

# ═══ ラベル ═══
txt(Inches(0.6), Inches(6.35), Inches(7), Inches(0.2),
    "■ 実線 = 現在稼働中の接続    ■ 各色 = 機能と対応サービスの紐付き", sz=8, color=GRAY)


# ════════════════════════════════════════
# 右半分: 将来の拡張ポイント
# ════════════════════════════════════════
x_r = Inches(8.0)

rrect(x_r, Inches(1.2), Inches(5.0), Inches(0.5), DARK)
txt(x_r + Inches(0.2), Inches(1.25), Inches(4.6), Inches(0.4),
    "将来の拡張ポイント & 概算費用", sz=15, color=WHITE, bold=True)

# 拡張カード
extensions = [
    ("v1.0", "得点リモコン",
     "別デバイスからスコア操作\n撮影者は撮影に集中",
     "WebSocket / Supabase Realtime",
     "約30〜50万円", "2-3週間",
     PRIMARY, P_LIGHT),
    ("v1.5", "マルチカメラ手動切替",
     "複数スマホの映像を\n手動で切り替え表示",
     "LiveKit マルチトラック",
     "約80〜120万円", "1-2ヶ月",
     BLUE, B_LIGHT),
    ("v2.0", "AI自動アングル切替",
     "AIがベストアングルを\n判断して自動切替",
     "AI推論サーバー（GPU）",
     "約200〜350万円", "3-6ヶ月",
     ORANGE, O_LIGHT),
    ("v2.5", "AIハイライト自動生成",
     "試合のハイライトを\nAIが自動編集・生成",
     "動画解析AI + FFmpeg",
     "約150〜250万円", "2-4ヶ月",
     GREEN, G_LIGHT),
    ("v3.0", "SNS自動投稿",
     "ハイライトをInstagram\nTikTokに自動投稿",
     "各SNS API連携",
     "約50〜80万円", "2-3週間",
     PURPLE, P2_LIGHT),
]

for i, (ver, title, desc, tech, cost, period, ac, bg) in enumerate(extensions):
    y = Inches(1.85 + i * 1.05)
    rrect(x_r, y, Inches(5.0), Inches(0.9), bg, ac)
    # バージョンバッジ
    rrect(x_r + Inches(0.1), y + Inches(0.1), Inches(0.55), Inches(0.3), ac)
    txt(x_r + Inches(0.1), y + Inches(0.1), Inches(0.55), Inches(0.3),
        ver, sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # タイトル
    txt(x_r + Inches(0.75), y + Inches(0.08), Inches(1.5), Inches(0.25),
        title, sz=12, color=DARK, bold=True)
    # 説明
    txt(x_r + Inches(0.75), y + Inches(0.35), Inches(1.5), Inches(0.5),
        desc, sz=8, color=SUB)
    # 技術
    txt(x_r + Inches(2.4), y + Inches(0.08), Inches(1.3), Inches(0.2),
        "技術:", sz=7, color=GRAY)
    txt(x_r + Inches(2.4), y + Inches(0.25), Inches(1.3), Inches(0.3),
        tech, sz=8, color=TEXT)
    # 費用・期間
    rrect(x_r + Inches(3.8), y + Inches(0.1), Inches(1.05), Inches(0.55), WHITE, BORDER)
    txt(x_r + Inches(3.85), y + Inches(0.12), Inches(0.95), Inches(0.25),
        cost, sz=10, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    txt(x_r + Inches(3.85), y + Inches(0.4), Inches(0.95), Inches(0.2),
        period, sz=8, color=SUB, align=PP_ALIGN.CENTER)

# 合計
rrect(x_r, Inches(7.05), Inches(5.0), Inches(0.35), P_LIGHT, PRIMARY)
txt(x_r + Inches(0.2), Inches(7.08), Inches(3.5), Inches(0.3),
    "全拡張の概算合計: 約510〜850万円（段階的に投資）", sz=11, color=PRIMARY, bold=True)

# フッター
txt(Inches(0.6), Inches(7.2), Inches(7), Inches(0.2),
    "LIVE SPOtCH  |  LIN-NAH株式会社  |  2026年4月", sz=8, color=GRAY)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LIVE_SPOtCH_構造チャート.pptx")
prs.save(out)
print(f"保存完了: {out}")
