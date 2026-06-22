#!/usr/bin/env python3
"""LIVE SPOtCH 事業計画書 v3 — 2026/05/29 改訂版

変更点 (v2 → v3):
- 表紙日付を 2026/05/29 に・版を v3 に
- スケジュールを新フローに更新: 5/29 本依頼書作成 → 6/1 税理士確認 → 6 月上旬 銀行ご提出 → 6 月中旬 着金・正式ローンチ + 広告開始
- 5 月完璧化進捗を反映: OAuth 通過 (5/13) ✅・YouTube 埋め込み許可・モバイル遷移体感改善・依存ライブラリ最新化 (5/28-29) を完了として記載
- PR マージ件数を 5/10 時点 126 件 → 5/29 時点 127 件 に更新
- NEW: 詳細市場推計スライドを追加 (文科省・スポーツ庁・中体連・高体連・スポ少データ + TAM-SAM-SOM 視覚化)
- NEW: ローンチ前の実課金トラクションスライドを追加 (5/21 確認: 外部有料 3 件・MRR ¥1,300/月)
- まとめスライドを最新数値で更新
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FONT = "Meiryo UI"

# === カラーパレット（ホワイトベース） ===
BG       = RGBColor(0xFF, 0xFF, 0xFF)
CARD     = RGBColor(0xF8, 0xF8, 0xFA)
CARD2    = RGBColor(0xF0, 0xF0, 0xF5)
TEXT     = RGBColor(0x1A, 0x1A, 0x2E)
SUB      = RGBColor(0x66, 0x66, 0x80)
BORDER   = RGBColor(0xE0, 0xE0, 0xE8)
PRIMARY  = RGBColor(0xE6, 0x39, 0x46)
P_LIGHT  = RGBColor(0xFD, 0xEB, 0xED)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLUE     = RGBColor(0x25, 0x63, 0xEB)
B_LIGHT  = RGBColor(0xEB, 0xF0, 0xFF)
ORANGE   = RGBColor(0xEA, 0x88, 0x0B)
O_LIGHT  = RGBColor(0xFF, 0xF3, 0xE0)
GREEN    = RGBColor(0x16, 0xA3, 0x4A)
G_LIGHT  = RGBColor(0xE8, 0xF5, 0xEB)
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)
P2_LIGHT = RGBColor(0xF3, 0xEE, 0xFF)
DARK     = RGBColor(0x0F, 0x17, 0x2A)
GRAY400  = RGBColor(0x9C, 0xA3, 0xAF)
RED_DARK = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# === ヘルパー関数 ===

def set_bg(sl, color=BG):
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = color

def rect(sl, l, t, w, h, color, lc=None):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if lc: s.line.color.rgb = lc; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def rrect(sl, l, t, w, h, color, lc=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if lc: s.line.color.rgb = lc; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def circle(sl, l, t, sz, color):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Pt(sz), Pt(sz))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def arrow_r(sl, l, t, w, h, color):
    s = sl.shapes.add_shape(MSO_SHAPE.NOTCHED_RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def txt(sl, l, t, w, h, text, sz=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = sl.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = FONT
    p.alignment = align; p.space_after = Pt(0); p.space_before = Pt(0)
    return box

def mtxt(sl, l, t, w, h, lines):
    box = sl.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, ld in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ld[0]; p.font.size = Pt(ld[1] if len(ld)>1 else 14)
        p.font.color.rgb = ld[2] if len(ld)>2 else TEXT
        p.font.bold = ld[3] if len(ld)>3 else False; p.font.name = FONT
        p.alignment = ld[4] if len(ld)>4 else PP_ALIGN.LEFT; p.space_after = Pt(4)
    return box

def hdr(sl, title, subtitle=""):
    set_bg(sl)
    rect(sl, Inches(0), Inches(0), Pt(5), Inches(7.5), PRIMARY)
    txt(sl, Inches(0.9), Inches(0.4), Inches(10), Inches(0.6), title, sz=30, color=DARK, bold=True)
    if subtitle:
        txt(sl, Inches(0.9), Inches(0.95), Inches(10), Inches(0.3), subtitle, sz=13, color=SUB)
    rect(sl, Inches(0.9), Inches(1.35), Inches(1.2), Pt(3), PRIMARY)

def pnum(sl, n):
    txt(sl, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.3), str(n), sz=10, color=GRAY400, align=PP_ALIGN.RIGHT)

def add_table(sl, l, t, w, rows, fsz=11):
    nr, nc = len(rows), len(rows[0])
    ts = sl.shapes.add_table(nr, nc, l, t, w, Inches(0.4*nr)); tbl = ts.table
    for r, rd in enumerate(rows):
        for c, ct in enumerate(rd):
            cell = tbl.cell(r,c); cell.text = str(ct); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(fsz); p.font.name = FONT; p.space_after = Pt(1)
                if r == 0:
                    p.font.bold = True; p.font.color.rgb = WHITE
                    cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
                else:
                    p.font.color.rgb = TEXT; cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE if r%2==1 else CARD
    return ts


# ════════════════════════════════════════════════
# 1. 表紙
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK)
rect(s, Inches(0), Inches(0), Pt(5), Inches(7.5), PRIMARY)

txt(s, Inches(1.4), Inches(1.5), Inches(10), Inches(0.6),
    "LIVE SPOtCH", sz=64, color=PRIMARY, bold=True)
txt(s, Inches(1.4), Inches(2.5), Inches(10), Inches(0.5),
    "誰もがスポーツ中継のカメラマン。", sz=28, color=WHITE, bold=True)
txt(s, Inches(1.4), Inches(3.1), Inches(10), Inches(0.5),
    "手元のスマホが機材になる。その一瞬の感動を手のひらに。", sz=22, color=RGBColor(0xCC, 0xCC, 0xDD))
rect(s, Inches(1.4), Inches(3.9), Inches(3), Pt(2), PRIMARY)
mtxt(s, Inches(1.4), Inches(4.3), Inches(8), Inches(2.5), [
    ("ローカルスポーツ ライブ配信プラットフォーム", 16, RGBColor(0x99, 0x99, 0xAA)),
    ("事業計画書 v3", 20, WHITE, True),
    ("", 8),
    ("LIN-NAH株式会社", 18, WHITE, True),
    ("2026年5月29日（v3 改訂版・市場規模推計+実課金トラクション増強）", 13, RGBColor(0x88, 0x88, 0x99)),
    ("", 8),
    ("https://live-spotch.com", 12, PRIMARY),
])


# ════════════════════════════════════════════════
# 2. ビジョン
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "ビジョン", "Vision"); pnum(s, 2)

# メインビジョン
rrect(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.2), P_LIGHT, PRIMARY)
txt(s, Inches(1.3), Inches(1.95), Inches(10.8), Inches(0.4),
    "「すべての子どもの試合が、どこからでも見られる世界をつくる」",
    sz=22, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.3), Inches(2.45), Inches(10.8), Inches(0.4),
    "スポーツブルが配信しない、地域大会・部活・少年団の試合を保護者の手で配信できる",
    sz=14, color=SUB, align=PP_ALIGN.CENTER)

# 3つの価値
values = [
    ("スマホ1台でOK", "プロカメラ機材は不要。\n保護者のスマホで\nTV中継品質の配信を実現。", PRIMARY, P_LIGHT),
    ("加速度的な拡散", "部活動・少年団は口コミが命。\n1チームの導入が対戦相手→\n大会→地域へとドミノ拡散。", BLUE, B_LIGHT),
    ("YouTube Live 連携", "チームプランは YouTube Live\n同時配信 + アーカイブ自動保存\n— 視聴も拡散も両立。", GREEN, G_LIGHT),
]
for i, (title, body, ac, bg_c) in enumerate(values):
    x = Inches(0.9 + i * 4.0)
    rrect(s, x, Inches(3.4), Inches(3.6), Inches(3.6), bg_c, ac)
    rect(s, x, Inches(3.4), Inches(3.6), Pt(4), ac)
    txt(s, x+Inches(0.3), Inches(3.7), Inches(3.0), Inches(0.3), title, sz=18, color=ac, bold=True)
    txt(s, x+Inches(0.3), Inches(4.2), Inches(3.0), Inches(2.5), body, sz=14, color=TEXT)


# ════════════════════════════════════════════════
# 3. 市場の課題
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "市場の課題", "Why Now? — 既存サービスがカバーしないローカルスポーツ"); pnum(s, 3)

# 左: 現状の課題
rrect(s, Inches(0.9), Inches(1.8), Inches(5.2), Inches(5.2), CARD, BORDER)
rect(s, Inches(0.9), Inches(1.8), Inches(5.2), Pt(4), RED_DARK)
txt(s, Inches(1.3), Inches(2.0), Inches(4.4), Inches(0.3),
    "現在のスポーツ配信市場の空白", sz=16, color=RED_DARK, bold=True)
problems = [
    "スポーツブル等は有名校・上位校の試合のみ",
    "地域大会・少年団の試合は一切配信されない",
    "YouTube Live ではスコアが表示できない",
    "BAND (NAVER) はチーム連絡網が中心、配信補助",
    "得点板を映すと映像が遮られる",
    "試合のアーカイブが散逸・管理できない",
]
for i, p in enumerate(problems):
    y = Inches(2.6 + i * 0.65)
    rrect(s, Inches(1.2), y, Inches(4.6), Inches(0.5), RGBColor(0xFF, 0xF0, 0xF0), RGBColor(0xFC, 0xD5, 0xD5))
    txt(s, Inches(1.4), y+Inches(0.06), Inches(4.2), Inches(0.4), f"✕  {p}", sz=13, color=RED_DARK)

# 矢印
arrow_r(s, Inches(6.3), Inches(3.8), Inches(1.0), Inches(0.5), PRIMARY)

# 右: 解決策
rrect(s, Inches(7.5), Inches(1.8), Inches(5.2), Inches(5.2), G_LIGHT, GREEN)
rect(s, Inches(7.5), Inches(1.8), Inches(5.2), Pt(4), GREEN)
txt(s, Inches(7.9), Inches(2.0), Inches(4.4), Inches(0.3),
    "LIVE SPOtCH の解決策", sz=16, color=GREEN, bold=True)
solutions = [
    ("スコアオーバーレイ", "TV中継のようにスコアを常時表示"),
    ("YouTube Live 同時配信", "チームプランで自動 push + アーカイブ"),
    ("Web版（PWA）", "URLタップで即視聴、DL不要"),
    ("LINE共有", "ワンタップで家族・チームに共有"),
    ("チーム管理", "メンバー招待・権限管理を一元化"),
    ("Made for Kids 対応", "13歳未満視聴対応で安全配信"),
]
for i, (title, desc) in enumerate(solutions):
    y = Inches(2.6 + i * 0.65)
    rrect(s, Inches(7.8), y, Inches(4.6), Inches(0.5), RGBColor(0xE8, 0xF8, 0xEB), RGBColor(0xBB, 0xE8, 0xC4))
    txt(s, Inches(8.0), y+Inches(0.06), Inches(2.0), Inches(0.4), f"✓  {title}", sz=12, color=GREEN, bold=True)
    txt(s, Inches(10.2), y+Inches(0.06), Inches(2.2), Inches(0.4), desc, sz=11, color=SUB)


# ════════════════════════════════════════════════
# 4. バイラル拡散の仕組み
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "なぜ部活動か — 加速度的な拡散モデル", "Viral Growth Strategy"); pnum(s, 4)

# ドミノフロー
steps = [
    ("Step 1", "保護者1人が配信", "スマホ1台で\n我が子の試合を配信\n→ チーム内30-50人が視聴", PRIMARY),
    ("Step 2", "対戦相手が発見", "「うちもやりたい！」\n→ 相手チームの保護者も\n配信者に転換", BLUE),
    ("Step 3", "大会全体に波及", "大会主催者が推奨\n→ 全参加チームが利用\n→ 1大会で100-200人獲得", GREEN),
    ("Step 4", "地域のスタンダード", "サッカー→野球→バスケ\n→ 複数スポーツに横展開\n→ 地域のインフラに", ORANGE),
]
rect(s, Inches(1.5), Inches(3.0), Inches(10.5), Pt(3), BORDER)
for i, (step, title, desc, ac) in enumerate(steps):
    x = Inches(0.9 + i * 3.1)
    c = circle(s, x + Inches(1.15), Inches(2.0), 45, ac)
    txt(s, x + Inches(1.05), Inches(2.03), Inches(0.5), Inches(0.4), str(i+1), sz=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    rrect(s, x, Inches(3.6), Inches(2.8), Inches(3.2), CARD, BORDER)
    rect(s, x, Inches(3.6), Inches(2.8), Pt(4), ac)
    txt(s, x+Inches(0.25), Inches(3.8), Inches(2.3), Inches(0.25), step, sz=11, color=ac, bold=True)
    txt(s, x+Inches(0.25), Inches(4.1), Inches(2.3), Inches(0.3), title, sz=16, color=DARK, bold=True)
    txt(s, x+Inches(0.25), Inches(4.6), Inches(2.3), Inches(2), desc, sz=13, color=SUB)
    if i < 3:
        txt(s, x + Inches(2.75), Inches(2.8), Inches(0.5), Inches(0.4), "→", sz=22, color=GRAY400, bold=True)

# 強調
rrect(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4), P_LIGHT, PRIMARY)
txt(s, Inches(1.3), Inches(7.02), Inches(10.8), Inches(0.35),
    "部活動は保護者コミュニティが密接 → 口コミの拡散速度が圧倒的に速い",
    sz=14, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# 5. 市場規模 + 競合
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "市場規模 + 競合分析", "Market Size & Competitive Landscape"); pnum(s, 5)

rrect(s, Inches(0.9), Inches(1.8), Inches(5.4), Inches(5.2), CARD, BORDER)
txt(s, Inches(1.3), Inches(1.95), Inches(4.6), Inches(0.3), "ターゲット市場規模", sz=16, color=DARK, bold=True)

bars = [("中学校 部活動", 270, 510), ("高校 部活動", 180, 510), ("スポーツ少年団", 60, 510)]
for i, (label, val, max_val) in enumerate(bars):
    y = Inches(2.6 + i * 1.0)
    txt(s, Inches(1.3), y, Inches(2.5), Inches(0.3), label, sz=13, color=SUB)
    txt(s, Inches(4.2), y, Inches(1.5), Inches(0.3), f"約{val}万人", sz=13, color=DARK, bold=True, align=PP_ALIGN.RIGHT)
    bar_w = Inches(4.0) * val / max_val
    rrect(s, Inches(1.3), y + Inches(0.35), Inches(4.3), Inches(0.25), CARD2)
    rrect(s, Inches(1.3), y + Inches(0.35), bar_w, Inches(0.25), PRIMARY)

rect(s, Inches(1.3), Inches(5.5), Inches(4.3), Pt(1), BORDER)
txt(s, Inches(1.3), Inches(5.7), Inches(4.3), Inches(0.3), "保護者・家族を含む対象層", sz=12, color=SUB)
txt(s, Inches(1.3), Inches(6.1), Inches(4.3), Inches(0.5),
    "1,000〜1,500万人", sz=30, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# 右: 競合ポジショニング (BAND を追加して 5 行)
rrect(s, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.2), CARD, BORDER)
txt(s, Inches(7.2), Inches(1.95), Inches(5), Inches(0.3), "競合との差別化", sz=16, color=DARK, bold=True)
comps = [
    ("スポーツブル", "プロ・大学スポーツ", "ローカル未対応", RED_DARK),
    ("YouTube Live", "汎用配信", "スコア・連携なし", RED_DARK),
    ("BAND (NAVER)", "チーム連絡網", "配信は補助機能", RED_DARK),
    ("DAZN", "プロ月額", "アマチュア未対応", RED_DARK),
    ("LIVE SPOtCH", "ローカルスポーツ専用", "唯一の選択肢", GREEN),
]
for i, (name, feat, gap, color) in enumerate(comps):
    y = Inches(2.6 + i * 0.85)
    is_us = name == "LIVE SPOtCH"
    bg_c = G_LIGHT if is_us else WHITE
    lc = GREEN if is_us else BORDER
    rrect(s, Inches(7.2), y, Inches(5), Inches(0.65), bg_c, lc)
    txt(s, Inches(7.5), y+Inches(0.1), Inches(1.8), Inches(0.3), name, sz=13, color=DARK, bold=is_us)
    txt(s, Inches(9.3), y+Inches(0.1), Inches(1.5), Inches(0.3), feat, sz=11, color=SUB)
    txt(s, Inches(10.8), y+Inches(0.1), Inches(1.3), Inches(0.3), gap, sz=11, color=color, bold=True)

rrect(s, Inches(6.8), Inches(7.05), Inches(5.6), Inches(0.35), P_LIGHT, PRIMARY)
txt(s, Inches(7.0), Inches(7.07), Inches(5.2), Inches(0.3),
    "国内に直接競合なし = ブルーオーシャン", sz=12, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# 6. NEW (v3): 詳細市場推計 — 公開統計データに基づく内訳
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "詳細市場推計", "Detailed Market Sizing — 公開統計データに基づく推計"); pnum(s, 6)

# 左: 公開データに基づくセグメント別人数
rrect(s, Inches(0.9), Inches(1.8), Inches(6.0), Inches(4.9), CARD, BORDER)
txt(s, Inches(1.2), Inches(1.95), Inches(5.4), Inches(0.3),
    "セグメント別の競技人口（公開統計）", sz=15, color=DARK, bold=True)

mkt_rows = [
    ("小学校 児童", "約 620 万人", "うちスポ少登録 約 60 万人", "文部科学省 学校基本調査"),
    ("中学校 運動部所属", "約 190 万人", "中学生の約 60% が運動部所属", "日本中学校体育連盟"),
    ("高校 運動部所属", "約 110 万人", "高校生の約 38% が運動部所属", "全国高等学校体育連盟"),
    ("スポーツ少年団 登録", "約 60 万人", "小学生中心の地域クラブ", "日本スポーツ少年団"),
]
for i, (seg, num, note, src) in enumerate(mkt_rows):
    y = Inches(2.5 + i * 0.95)
    rrect(s, Inches(1.2), y, Inches(5.4), Inches(0.8), WHITE, BORDER)
    txt(s, Inches(1.4), y+Inches(0.05), Inches(2.6), Inches(0.3), seg, sz=12, color=DARK, bold=True)
    txt(s, Inches(4.0), y+Inches(0.05), Inches(2.5), Inches(0.3), num, sz=14, color=PRIMARY, bold=True, align=PP_ALIGN.RIGHT)
    txt(s, Inches(1.4), y+Inches(0.4), Inches(5.0), Inches(0.2), note, sz=10, color=SUB)
    txt(s, Inches(1.4), y+Inches(0.58), Inches(5.0), Inches(0.2), f"出典: {src}", sz=9, color=GRAY400)

rrect(s, Inches(1.2), Inches(6.4), Inches(5.4), Inches(0.3), P_LIGHT, PRIMARY)
txt(s, Inches(1.4), Inches(6.42), Inches(5.0), Inches(0.25),
    "コア競技人口 合計  約 360-380 万人", sz=12, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# 右: TAM-SAM-SOM funnel
rrect(s, Inches(7.2), Inches(1.8), Inches(5.4), Inches(4.9), CARD, BORDER)
txt(s, Inches(7.5), Inches(1.95), Inches(4.8), Inches(0.3),
    "TAM / SAM / SOM 推計", sz=15, color=DARK, bold=True)

tams = [
    ("TAM", "全潜在視聴・配信層", "1,140-1,520 万人", "競技人口 380万 × 平均 3-4 名 (両親・祖父母・近親者)", PRIMARY, P_LIGHT, 5.0),
    ("SAM", "撮影可能保護者", "190-380 万人", "競技人口の保護者 1-2 名 (配信者候補)", BLUE, B_LIGHT, 3.8),
    ("SOM", "Year 3 達成目標", "13,000 人 / 5,144 万円", "SAM の約 0.5% 以下 (慎重な算定)", GREEN, G_LIGHT, 2.6),
]
for i, (label, title, num, note, ac, bg_c, w) in enumerate(tams):
    y = Inches(2.5 + i * 1.25)
    x = Inches(7.5 + (5.4 - w) / 2)
    rrect(s, x, y, Inches(w), Inches(1.05), bg_c, ac)
    rect(s, x, y, Pt(4), Inches(1.05), ac)
    txt(s, x+Inches(0.2), y+Inches(0.06), Inches(w-0.4), Inches(0.3), f"{label} — {title}", sz=12, color=ac, bold=True)
    txt(s, x+Inches(0.2), y+Inches(0.35), Inches(w-0.4), Inches(0.35), num, sz=18, color=DARK, bold=True)
    txt(s, x+Inches(0.2), y+Inches(0.75), Inches(w-0.4), Inches(0.25), note, sz=9, color=SUB)

# フッター: 出典まとめ + 補足
rrect(s, Inches(0.9), Inches(7.05), Inches(11.7), Inches(0.35), CARD2, BORDER)
txt(s, Inches(1.1), Inches(7.07), Inches(11.3), Inches(0.3),
    "※ 数値は文部科学省「学校基本調査 令和5年度」/ 中体連・高体連 登録者数 / 日本スポーツ少年団 登録状況 等の公開データに基づく当社推計",
    sz=10, color=SUB)


# ════════════════════════════════════════════════
# 7. ビジネスモデル v2 (4/30 確定)  ※ v3 でスライド番号 6→7 に繰り下げ
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "ビジネスモデル", "Pricing Plans v2 (2026/04/30 確定)"); pnum(s, 7)

plans = [
    ("無料プラン","¥0","視聴者向け",["ライブ視聴（無制限）","1ヶ月以内のアーカイブ視聴","試合スコアの閲覧"],SUB,CARD,BORDER, False),
    ("配信者プラン","¥300/月","個人保護者向け",["ライブ配信機能（無制限）","スコアボード常時表示","LINE共有・QRコード","限定公開URL","初回10分間無料お試し","※ アーカイブなし(ライブ専用)"],PRIMARY,P_LIGHT,PRIMARY, False),
    ("チームプラン","¥500/月","チーム代表向け",["配信者プランの全機能","YouTube Live 同時配信","YouTubeアーカイブ自動保存","チーム管理・メンバー招待","スケジュール管理","AIハイライト (開発中)"],ORANGE,O_LIGHT,ORANGE, True),
]
for i, (name,price,tgt,feats,ac,bg_c,bdr,is_recommended) in enumerate(plans):
    x = Inches(0.9+i*4.1); w = Inches(3.8)
    rrect(s, x, Inches(1.8), w, Inches(5.2), bg_c, bdr)
    rect(s, x, Inches(1.8), w, Pt(4), ac)
    if is_recommended:
        rrect(s, x+Inches(2.2), Inches(1.6), Inches(1.4), Inches(0.35), ORANGE)
        txt(s, x+Inches(2.2), Inches(1.6), Inches(1.4), Inches(0.35), "推奨", sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.3), Inches(2.1), Inches(3.2), Inches(0.3), name, sz=16, color=ac, bold=True)
    txt(s, x+Inches(0.3), Inches(2.5), Inches(3.2), Inches(0.5), price, sz=38, color=DARK, bold=True)
    txt(s, x+Inches(0.3), Inches(3.15), Inches(3.2), Inches(0.3), tgt, sz=12, color=SUB)
    rect(s, x+Inches(0.3), Inches(3.5), w-Inches(0.6), Pt(1), BORDER)
    for j, f in enumerate(feats):
        txt(s, x+Inches(0.3), Inches(3.7+j*0.42), Inches(3.4), Inches(0.3), f"  ✓  {f}", sz=12, color=TEXT)

# 補足: 1ヶ月後アーカイブ課金
rrect(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.35), B_LIGHT, BLUE)
txt(s, Inches(1.3), Inches(7.02), Inches(10.8), Inches(0.3),
    "アーカイブ課金 ¥100/視聴: 1ヶ月経過後のアーカイブはチームプラン配信のみ視聴可能 (容量ベース従量課金)",
    sz=11, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# 7. 3年間の収支計画
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "3年間の収支計画", "Financial Projection"); pnum(s, 8)

rows = [["","Year 1（立ち上げ期）","Year 2（成長期）","Year 3（拡大期）"],
    ["登録ユーザー（期末）","8,000人","35,000人","100,000人"],
    ["有料会員（期末）","1,000人","4,550人","13,000人"],
    ["売上合計","153万円","1,291万円","5,144万円"],
    ["経費合計","503万円","634万円","1,479万円"],
    ["営業損益","▲350万円","+657万円","+3,665万円"],
    ["累計損益","▲350万円","+307万円","+3,972万円"]]
add_table(s, Inches(0.9), Inches(1.8), Inches(11.5), rows, fsz=14)

rrect(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.6), P_LIGHT, PRIMARY)
txt(s, Inches(1.3), Inches(5.68), Inches(10.8), Inches(0.45),
    "Year 2で単月黒字化 → 借入金700万円はYear 2〜3で完済見込み",
    sz=16, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# 8. 資金計画 v2 (700万配分更新)
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "資金計画", "700万円の使途内訳 (v2 改訂)"); pnum(s, 9)

rrect(s, Inches(0.9), Inches(1.8), Inches(3.5), Inches(1.2), P_LIGHT, PRIMARY)
mtxt(s, Inches(1.3), Inches(1.85), Inches(2.8), Inches(1), [
    ("借入希望額", 12, SUB), ("700万円", 42, PRIMARY, True)])
mtxt(s, Inches(0.9), Inches(3.05), Inches(3.5), Inches(0.4), [
    ("公庫350万 + 埼玉りそな350万", 11, SUB),
    ("協調融資による分散リスク", 10, SUB),
])

items = [
    ("マーケティング費","240万円","Meta/Google広告 / PRTIMES / SNS / 販促物",34, PRIMARY),
    ("開発費","210万円","LiveKit Cloud Ship + TCP化中継server + AI開発支援",30, BLUE),
    ("運転資金","170万円","LIN-NAH既存事業との並行運営 (6ヶ月分)",24, ORANGE),
    ("予備費","50万円","OAuth審査追加対応 / TCP化想定外コスト / 法務",7, PURPLE),
    ("インフラ費","30万円","YouTube API / Egress帯域 / Vercel Pro",5, GREEN),
]

for i, (name,amt,desc,pct,color) in enumerate(items):
    y = Inches(3.3 + i * 0.85)
    txt(s, Inches(0.9), y, Inches(2.2), Inches(0.3), name, sz=14, color=DARK, bold=True)
    txt(s, Inches(3.3), y, Inches(1.5), Inches(0.3), amt, sz=14, color=color, bold=True, align=PP_ALIGN.RIGHT)
    max_w = Inches(7.5)
    bar_w = max_w * pct / 40
    rrect(s, Inches(5.0), y + Inches(0.05), max_w, Inches(0.25), CARD2)
    rrect(s, Inches(5.0), y + Inches(0.05), bar_w, Inches(0.25), color)
    txt(s, Inches(5.0) + bar_w + Inches(0.1), y, Inches(1), Inches(0.3), f"{pct}%", sz=11, color=color)
    txt(s, Inches(5.0), y + Inches(0.35), Inches(7.5), Inches(0.25), desc, sz=10, color=SUB)


# ════════════════════════════════════════════════
# 9. NEW: 5月中完璧化 + 6/1 ローンチ計画
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "完璧化進捗 → 6月中旬 正式ローンチ", "Phase 0: Pre-Launch Polish 完了状況 + 6月実行プラン"); pnum(s, 10)

# 左: 5月完璧化の3本柱 (進捗反映)
rrect(s, Inches(0.9), Inches(1.8), Inches(6.0), Inches(5.4), CARD, BORDER)
txt(s, Inches(1.3), Inches(1.95), Inches(5.4), Inches(0.3),
    "5月完璧化 — 3本柱の進捗", sz=16, color=DARK, bold=True)

pillars = [
    ("A. Google OAuth 認定通過 ✅ 完了", "2026-05-13 正式通過済。3 scope (readonly/upload/youtube) すべて承認。「未確認のアプリ」警告画面消滅 + 100枠制限解除", GREEN),
    ("C. 本番運用上の細部仕上げ ✅ 完了", "5/28-29 に YouTube 埋め込み許可・モバイル遷移体感改善・依存ライブラリ 13 種一括最新化・「アプリを開く」体感空白解消を実施。本番品質を引き上げ", GREEN),
    ("B. 配信プロトコル TCP 化 → ローンチ後継続改善", "MediaRecorder + 中継 server で RTMP 化し 4G 環境下の高画質化。月+$5-20。ローンチ品質に必須ではないため継続改善項目に整理", BLUE),
]
for i, (title, body, ac) in enumerate(pillars):
    y = Inches(2.5 + i * 1.5)
    rrect(s, Inches(1.2), y, Inches(5.4), Inches(1.4), WHITE, ac)
    rect(s, Inches(1.2), y, Pt(4), Inches(1.4), ac)
    txt(s, Inches(1.4), y+Inches(0.1), Inches(5.0), Inches(0.3), title, sz=13, color=ac, bold=True)
    txt(s, Inches(1.4), y+Inches(0.5), Inches(5.0), Inches(0.9), body, sz=11, color=TEXT)

# 右: タイムライン (v3 新スケジュール)
rrect(s, Inches(7.2), Inches(1.8), Inches(5.4), Inches(5.4), P_LIGHT, PRIMARY)
txt(s, Inches(7.5), Inches(1.95), Inches(4.8), Inches(0.3),
    "ローンチタイムライン (v3)", sz=16, color=PRIMARY, bold=True)

milestones = [
    ("5/10", "本番稼働・ブロックシード大会 53 分配信完走", DARK),
    ("5/13", "✅ Google OAuth 検証通過 (3 scope 承認)", GREEN),
    ("5/21", "✅ 無施策で外部有料 3 件・MRR ¥1,300/月 確認", GREEN),
    ("5/28-29", "✅ 視聴・モバイル・依存 一括改善 (127 PR)", GREEN),
    ("5/29", "本依頼書作成 (本日)", PRIMARY),
    ("6/1", "★ 顧問税理士確認・最終化", PRIMARY),
    ("6月上旬", "★ 公庫 + 埼玉りそな ご提出", PRIMARY),
    ("6月中旬", "★ 着金 + 正式ローンチ + 広告開始", PRIMARY),
    ("8月末", "KPI ゲート 1: 有料 5 名", GREEN),
    ("10月初", "KPI ゲート 2: 有料 20 名", GREEN),
]
for i, (date, body, color) in enumerate(milestones):
    y = Inches(2.5 + i * 0.5)
    is_star = body.startswith("★")
    txt(s, Inches(7.5), y, Inches(0.9), Inches(0.35), date, sz=12, color=color, bold=is_star)
    txt(s, Inches(8.5), y, Inches(4.0), Inches(0.35), body, sz=11, color=color, bold=is_star)


# ════════════════════════════════════════════════
# 11. NEW (v3): ローンチ前の実課金トラクション
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "ローンチ前の実課金トラクション", "Pre-Launch Paid Traction — 無施策で外部有料課金 3 件"); pnum(s, 11)

# 上部: MRR + 強調メッセージ
rrect(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.4), G_LIGHT, GREEN)
rect(s, Inches(0.9), Inches(1.8), Pt(4), Inches(1.4), GREEN)
txt(s, Inches(1.3), Inches(1.95), Inches(11), Inches(0.4),
    "正式ローンチ前・広告ゼロ・営業ゼロの段階で、外部から有料課金 3 件が自然発生",
    sz=15, color=GREEN, bold=True)
mtxt(s, Inches(1.3), Inches(2.4), Inches(11), Inches(0.8), [
    ("実課金 MRR", 11, SUB),
    ("¥1,300 / 月  (チーム¥500 × 2 + 配信者¥300 × 1)", 22, DARK, True),
])

# 下部: 3件の内訳テーブル
traction_rows = [
    ["カテゴリ (匿名化)", "プラン", "課金開始日", "実利用"],
    ["少年サッカーアカデミー (U8)", "チーム ¥500/月", "2026-05-18", "当日に配信 3 回"],
    ["スポーツチーム", "チーム ¥500/月", "2026-05-15", "2 日で配信 4 回"],
    ["個人保護者", "配信者 ¥300/月", "2026-05-08", "当日に配信 4 回"],
]
add_table(s, Inches(0.9), Inches(3.4), Inches(11.5), traction_rows, fsz=13)

# 下部: 強調メッセージ
rrect(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.3), P_LIGHT, PRIMARY)
txt(s, Inches(1.3), Inches(5.85), Inches(11), Inches(0.3),
    "「作ったが誰も使わない」リスクの低い、市場に求められた事業である実証", sz=14, color=PRIMARY, bold=True)
points_traction = [
    "3 件とも Stripe 本番決済での課金が継続中・全員が実際に試合配信に利用（計 11 配信）",
    "うち 2 件は地域の子どもスポーツチーム — ターゲット層そのものが自発的に発見・課金",
    "5/21 オーナー確認時点。個人情報保護のため氏名・チーム名は匿名化、Stripe・Supabase 記録にて実在確認可能",
]
for i, pt in enumerate(points_traction):
    txt(s, Inches(1.5), Inches(6.2+i*0.32), Inches(10.8), Inches(0.3), f"  ✓  {pt}", sz=11, color=TEXT)


# ════════════════════════════════════════════════
# 12. 開発実績 (5/29時点)  ※ v3 で 10→12 に繰り下げ
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "開発実績", "本番稼働中 — 直近50日で127件PRマージ"); pnum(s, 12)

# 進捗バー: 100%
rrect(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.5), CARD2, BORDER)
rrect(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.5), GREEN)
txt(s, Inches(0.9), Inches(1.83), Inches(11.5), Inches(0.45), "本番稼働 100%（基盤完成 / 細部仕上げ中）", sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 左: PR内訳
rrect(s, Inches(0.9), Inches(2.6), Inches(7.5), Inches(4.4), CARD, BORDER)
txt(s, Inches(1.3), Inches(2.7), Inches(6.5), Inches(0.3), "✅ 直近50日 (4/9-5/29) PR 内訳", sz=14, color=GREEN, bold=True)
pr_items = [
    ("インフラ・セキュリティ", "18 件"),
    ("配信品質改善 (A/V drift / カメラ / 音声)", "28 件"),
    ("YouTube連携 (OAuth / Live同時配信 / 埋め込み許可)", "19 件"),
    ("決済・課金 (Stripe)", "8 件"),
    ("UI/UX 改善 (モバイル遷移体感改善含む)", "31 件"),
    ("運用堅牢化 (ストレイ対策・視聴経路最適化)", "6 件"),
    ("その他 (依存更新・ドキュメント)", "17 件"),
    ("【合計】", "127 件"),
]
for i, (name, count) in enumerate(pr_items):
    y = Inches(3.15 + i * 0.42)
    is_total = "合計" in name
    txt(s, Inches(1.3), y, Inches(5.3), Inches(0.3), name, sz=12, color=DARK if is_total else TEXT, bold=is_total)
    txt(s, Inches(6.6), y, Inches(1.6), Inches(0.3), count, sz=12, color=PRIMARY if is_total else SUB, bold=is_total, align=PP_ALIGN.RIGHT)

# 右: 残タスク
rrect(s, Inches(8.7), Inches(2.6), Inches(3.7), Inches(2.5), O_LIGHT, ORANGE)
txt(s, Inches(9.0), Inches(2.7), Inches(3.1), Inches(0.3), "5月完璧化 進捗", sz=14, color=ORANGE, bold=True)
remaining = [
    ("Google OAuth 認定通過", "✅ 5/13"),
    ("細部 UX 仕上げ", "✅ 5/28-29"),
    ("配信 TCP 化 (B案)", "ローンチ後"),
]
for i, (task, period) in enumerate(remaining):
    y = Inches(3.2+i*0.6)
    txt(s, Inches(9.0), y, Inches(2.0), Inches(0.3), f"  {task}", sz=12, color=TEXT)
    txt(s, Inches(11.0), y, Inches(1.2), Inches(0.3), period, sz=11, color=SUB, align=PP_ALIGN.RIGHT)

# 本番URL
rrect(s, Inches(8.7), Inches(5.4), Inches(3.7), Inches(1.6), CARD, BORDER)
txt(s, Inches(9.0), Inches(5.5), Inches(3.1), Inches(0.25), "本番URL", sz=11, color=SUB)
txt(s, Inches(9.0), Inches(5.8), Inches(3.1), Inches(0.25), "live-spotch.com", sz=14, color=PRIMARY, bold=True)
txt(s, Inches(9.0), Inches(6.15), Inches(3.1), Inches(0.25), "本番稼働: 2026/4/19〜", sz=10, color=SUB)
txt(s, Inches(9.0), Inches(6.45), Inches(3.1), Inches(0.3), "正式ローンチ: 2026/6 中旬", sz=12, color=DARK, bold=True)

# 注釈
rrect(s, Inches(0.9), Inches(7.05), Inches(11.5), Inches(0.35), P_LIGHT, PRIMARY)
txt(s, Inches(1.3), Inches(7.07), Inches(10.8), Inches(0.3),
    "5/10 大会 53 分配信完走 / 5/13 OAuth 通過 / 5/21 外部有料 3 件確認 / 5/28-29 視聴・モバイル・依存 一括改善",
    sz=11, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# 11. 知的財産・参入障壁
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "知的財産・参入障壁", "IP & Competitive Moat"); pnum(s, 13)

items_ip = [
    ("商標登録", "「LIVE SPOtCH」の商標出願を予定。\nサービス名のブランド保護。", PRIMARY, P_LIGHT),
    ("特許出願", "スコアオーバーレイ×ライブ配信×\nスポーツルールエンジンの組み合わせ\n技術について特許出願を検討。", BLUE, B_LIGHT),
    ("先行者優位", "ネットワーク効果：チーム単位で普及\n→ 先に広がった側が勝つ。\nUGCの蓄積が参入障壁になる。", GREEN, G_LIGHT),
    ("AI技術蓄積", "将来のマルチカメラAI自動切替、\nAIハイライト自動生成など\n独自の技術蓄積を計画。", PURPLE, P2_LIGHT),
]
for i, (title, body, ac, bg_c) in enumerate(items_ip):
    x = Inches(0.9 + (i % 2) * 6.1)
    y = Inches(1.8 + (i // 2) * 2.8)
    rrect(s, x, y, Inches(5.6), Inches(2.4), bg_c, ac)
    rect(s, x, y, Inches(5.6), Pt(4), ac)
    txt(s, x+Inches(0.3), y+Inches(0.25), Inches(5), Inches(0.3), title, sz=18, color=ac, bold=True)
    txt(s, x+Inches(0.3), y+Inches(0.7), Inches(5), Inches(1.5), body, sz=14, color=TEXT)


# ════════════════════════════════════════════════
# 12. 開発コスト戦略
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); hdr(s, "開発コスト戦略", "AI活用による大幅なコスト削減"); pnum(s, 14)

rrect(s, Inches(0.9), Inches(1.8), Inches(5.2), Inches(2.2), CARD, BORDER)
rect(s, Inches(0.9), Inches(1.8), Inches(5.2), Pt(4), RED_DARK)
txt(s, Inches(1.3), Inches(2.0), Inches(4.4), Inches(0.3), "通常の開発会社に外注した場合", sz=14, color=RED_DARK, bold=True)
mtxt(s, Inches(1.3), Inches(2.5), Inches(4.4), Inches(1.2), [
    ("現在完成分の見積もり", 12, SUB),
    ("約 1,200 万円（税込）", 26, DARK, True)])

arrow_r(s, Inches(6.3), Inches(2.5), Inches(1.0), Inches(0.5), PRIMARY)

rrect(s, Inches(7.5), Inches(1.8), Inches(5.2), Inches(2.2), G_LIGHT, GREEN)
rect(s, Inches(7.5), Inches(1.8), Inches(5.2), Pt(4), GREEN)
txt(s, Inches(7.9), Inches(2.0), Inches(4.4), Inches(0.3), "AI開発支援活用の実績", sz=14, color=GREEN, bold=True)
mtxt(s, Inches(7.9), Inches(2.5), Inches(4.4), Inches(1.2), [
    ("実コスト", 12, SUB),
    ("ほぼ ¥0（自社開発）", 26, GREEN, True)])

rrect(s, Inches(0.9), Inches(4.3), Inches(11.8), Inches(2.8), CARD, BORDER)
txt(s, Inches(1.3), Inches(4.4), Inches(10), Inches(0.3),
    "なぜこれが可能なのか", sz=16, color=DARK, bold=True)
points = [
    "Claude 等の最新 AI 開発支援ツールにより、設計・実装・テスト・PR レビューまで自社で対応",
    "従来の外注開発に対し、90% 以上のコスト削減を実現",
    "5/29 時点で計 127 PR マージ、本番運用中（決済・チーム管理・YouTube連携すべて稼働中）",
    "5/10 試合実機で見えた課題を当日中に 5 PR で hotfix・5/28-29 にも視聴/モバイル一括改善 — 本番運用と並行した即応開発体制",
    "既存 LIN-NAH 事業で培ったデザイン・映像制作のノウハウもプロダクトに活用",
]
for i, pt in enumerate(points):
    txt(s, Inches(1.5), Inches(4.85+i*0.4), Inches(10.8), Inches(0.3), f"  ✓  {pt}", sz=12, color=TEXT)


# ════════════════════════════════════════════════
# 13. まとめ
# ════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK)
rect(s, Inches(0), Inches(0), Pt(5), Inches(7.5), PRIMARY)

txt(s, Inches(1.4), Inches(0.8), Inches(10), Inches(0.8), "LIVE SPOtCH", sz=56, color=PRIMARY, bold=True)
txt(s, Inches(1.4), Inches(1.7), Inches(10), Inches(0.5),
    "誰もがスポーツ中継のカメラマン。その一瞬の感動を手のひらに。", sz=22, color=WHITE)
rect(s, Inches(1.4), Inches(2.5), Inches(3), Pt(2), PRIMARY)

points_final = [
    "国内に直接競合がないブルーオーシャン市場（TAM 1,140-1,520 万人 / 公開統計に基づく推計）",
    "本番稼働 100% (live-spotch.com) — 5/29 時点で計 127 PR マージ、5/13 Google OAuth 認定通過",
    "ローンチ前・無施策で外部有料 3 件・MRR ¥1,300/月 を確認（市場ニーズの実証）",
    "部活動・少年団コミュニティの口コミで加速度的に拡散（バイラルモデル）",
    "UGCモデルで撮影コストゼロの全国展開が可能",
    "YouTube Live 同時配信 + アーカイブで使いやすさと拡散力を最大化",
    "Year 2 で単月黒字化、返済後も十分なキャッシュフロー確保",
    "6/1 税理士確認 → 6 月上旬 銀行ご提出 → 6 月中旬 着金・正式ローンチ + 広告開始",
]
for i, pt in enumerate(points_final):
    txt(s, Inches(1.4), Inches(3.0+i*0.48), Inches(10), Inches(0.4),
        f"  ✓  {pt}", sz=15, color=RGBColor(0xCC, 0xCC, 0xDD))

rrect(s, Inches(1.4), Inches(6.5), Inches(7.0), Inches(0.7), RGBColor(0x18, 0x18, 0x22), PRIMARY)
txt(s, Inches(1.7), Inches(6.58), Inches(6.5), Inches(0.5),
    "借入希望: 公庫350万 + 埼玉りそな350万 = 700万円", sz=16, color=PRIMARY, bold=True)
txt(s, Inches(8.7), Inches(6.6), Inches(4.0), Inches(0.4),
    "LIN-NAH 株式会社", sz=14, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
txt(s, Inches(8.7), Inches(6.95), Inches(4.0), Inches(0.3),
    "代表 原田 佑樹", sz=11, color=RGBColor(0xCC, 0xCC, 0xDD), align=PP_ALIGN.RIGHT)


# ── 保存 ──
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LIVE_SPOtCH_事業計画書.pptx")
prs.save(out)
print(f"保存完了: {out}")
print(f"")
print(f"=== 事業計画書 v3 / 2026-05-29 改訂 ===")
print(f"スライド数: {len(prs.slides)}  (v2: 13 → v3: 15)")
print(f"NEW スライド:")
print(f"  - #6  詳細市場推計 (公開統計データ + TAM/SAM/SOM)")
print(f"  - #11 ローンチ前の実課金トラクション (MRR ¥1,300/月)")
print(f"更新スライド:")
print(f"  - #1  表紙: v3・2026/05/29 改訂版")
print(f"  - #10 ローンチタイムライン (新スケジュール: 6/1 税理士 → 6 上旬 銀行 → 6 中旬 ローンチ)")
print(f"  - #12 開発実績 (5/29 時点 127 PR・5/13 OAuth 通過反映)")
print(f"  - #14 開発コスト戦略 (PR 数最新化)")
print(f"  - #15 まとめ (実課金実証・市場推計・新スケジュール反映)")
